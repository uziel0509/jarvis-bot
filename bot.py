"""
JARVIS 3.0 — Bot de Telegram con IA para estudiantes universitarios peruanos
Autor: uziel0509
Versión: 3.0.0
"""

import os
import json
import logging
import re
import asyncio
import sqlite3
import base64
import tempfile
from datetime import datetime, timedelta
from io import BytesIO
from pathlib import Path

from dotenv import load_dotenv
from groq import Groq
from telegram import Update, BotCommand
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    ContextTypes, filters, ConversationHandler
)
from apscheduler.schedulers.asyncio import AsyncIOScheduler
from apscheduler.triggers.date import DateTrigger

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.mathtext as mathtext

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    HRFlowable, Image as RLImage, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY, TA_RIGHT
from reportlab.pdfgen import canvas as rl_canvas

# ─────────────────────────────────────────────
# MÓDULOS JARVIS 3.1
# ─────────────────────────────────────────────
from modulos.finanzas import (
    registrar_gasto, registrar_ingreso, resumen_mensual,
    configurar_limite, agregar_pago_recurrente, check_pagos_hoy,
    get_pagos_hoy, prompt_interpretar_finanzas
)
from modulos.horario import (
    necesita_horario, mensaje_pedir_foto_horario, guardar_horario_extraido,
    registrar_examen, registrar_entrega, resumen_hoy,
    get_examenes_sin_analizar, mensaje_pedir_foto_examen,
    guardar_resultado_examen, PROMPT_EXTRAER_HORARIO, PROMPT_ANALIZAR_EXAMEN
)
from modulos.pre_render import procesar_output, elementos_a_texto_plano

# ─────────────────────────────────────────────
# CONFIGURACIÓN
# ─────────────────────────────────────────────
load_dotenv()

TELEGRAM_TOKEN  = os.getenv("TELEGRAM_TOKEN")
GROQ_API_KEY    = os.getenv("GROQ_API_KEY")
DB_PATH         = "/root/jarvis/jarvis.db"
HISTORIAL_DIR   = "/root/jarvis/historial"
PERFILES_DIR    = "/root/jarvis/perfiles"
ARCHIVOS_DIR    = "/root/jarvis/archivos"

for d in [HISTORIAL_DIR, PERFILES_DIR, ARCHIVOS_DIR]:
    Path(d).mkdir(parents=True, exist_ok=True)

# ─────────────────────────────────────────────
# MODELOS IA (optimizados por tarea)
# ─────────────────────────────────────────────
MODELO_CHAT              = "llama-3.3-70b-versatile"
MODELO_EJERCICIOS        = "openai/gpt-oss-120b"
MODELO_EJERCICIOS_RAPIDO = "openai/gpt-oss-120b"
MODELO_VISION            = "meta-llama/llama-4-scout-17b-16e-instruct"
MODELO_VOZ               = "whisper-large-v3-turbo"
MODELO_RAPIDO            = "llama-3.1-8b-instant"
MODELO_RESUMEN           = "openai/gpt-oss-20b"

# ─────────────────────────────────────────────
# LÍMITES POR PLAN
# ─────────────────────────────────────────────
PLAN_BASICO_LIMITE = 500
PLAN_FULL_LIMITE   = 99999

# ─────────────────────────────────────────────
# MEMORIA
# ─────────────────────────────────────────────
RECENT_MSGS   = 15
SUMMARY_EVERY = 20

# ─────────────────────────────────────────────
# ONBOARDING — estados de conversación
# ─────────────────────────────────────────────
OB_NOMBRE, OB_CARRERA, OB_CICLO, OB_UNIVERSIDAD = range(4)

# ─────────────────────────────────────────────
# ESTADO DE CONTEXTO DE FOTO — distingue tipo de imagen esperada
# ─────────────────────────────────────────────
# Se guarda en context.user_data["esperando_foto"]
FOTO_HORARIO  = "horario"
FOTO_EXAMEN   = "examen"
FOTO_EJERCICIO = "ejercicio"

# ─────────────────────────────────────────────
# LOGGING
# ─────────────────────────────────────────────
# Logging a consola + archivo
Path("/root/jarvis/logs").mkdir(parents=True, exist_ok=True)
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("/root/jarvis/logs/jarvis.log", encoding="utf-8")
    ]
)
logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────
# CLIENTES
# ─────────────────────────────────────────────
client    = Groq(api_key=GROQ_API_KEY)
scheduler = AsyncIOScheduler(timezone="America/Lima")

# ─────────────────────────────────────────────
# SYSTEM PROMPT BASE
# ─────────────────────────────────────────────
SYSTEM_PROMPT_BASE = """Eres JARVIS, el asistente de IA más avanzado para estudiantes universitarios peruanos.
Tu misión: resolver ejercicios académicos, organizar la vida estudiantil y ser el mejor socio de éxito.

PERSONALIDAD:
- Directo, confiado y motivador. Hablas como un compañero brillante, no como un robot corporativo.
- Adaptas tu tono: técnico con avanzados, didáctico con principiantes.
- Usas ejemplos del contexto peruano cuando puedes.
- Celebras los logros del alumno con entusiasmo genuino.
- Eres honesto: si algo está mal, lo dices directo, pero con respeto.

CAPACIDADES:
- Resuelves ejercicios de física, matemáticas, química, ingeniería con pasos detallados
- Generas PDFs profesionales con soluciones y fórmulas bien renderizadas
- Creas archivos Excel con fórmulas, gráficas y formato profesional
- Generas presentaciones PowerPoint con diseño visual
- Redactas documentos Word formateados
- Haces recordatorios inteligentes (entiendes fechas en lenguaje natural)
- Llevas la contabilidad personal del alumno
- Analizas imágenes con ejercicios y las resuelves paso a paso

REGLAS DE RESPUESTA OBLIGATORIAS:
- NUNCA uses LaTeX crudo en el chat: NO escribas \\frac, $$formula$$, \\alpha, \\int
- Escribe fórmulas en texto legible usando caracteres Unicode reales:
  * Subíndices Unicode: H₂O, CO₂, C₈H₁₈, m⁻¹, 10⁶, Nₐ, m³, mg·m⁻³
  * NO uses letras "n" como reemplazo: NUNCA escribas HnO, COn, mn1, 10n, Nn
  * Fracciones: (numerador / denominador)
  * Operadores: × · ÷ ± ≤ ≥ ≠ ≈ ∑ √ ∫ Δ
- En PDFs las fórmulas SÍ se renderizan visualmente
- Cuando resuelvas MÚLTIPLES ejercicios (Ejercicio 1, Ejercicio 2, etc.):
  * Cada ejercicio debe empezar con: ## Ejercicio N: [título descriptivo]
  * Los pasos de CADA ejercicio empiezan desde Paso 1 (no son continuos)
  * Ejemplo: Ejercicio 17 → Paso 1, Paso 2... RESULTADO. Ejercicio 18 → Paso 1, Paso 2... RESULTADO
  * Nunca continúes la numeración de pasos entre ejercicios distintos
- Cuando resuelvas UN solo ejercicio, inicia con: ## [título descriptivo del ejercicio]
- Responde siempre en español a menos que el alumno escriba en otro idioma
- NO uses frases robóticas como "¡Por supuesto!", "¡Claro que sí!", "Como asistente de IA"
- NO termines cada mensaje con "¿Hay algo más en lo que pueda ayudarte?" ni "¿En qué puedo ayudarte?"
- NO empieces saludos con preguntas genéricas de asistente
- Máximo 1 emoji por mensaje, úsalos solo cuando aporten algo
- Sé conciso: si la respuesta es corta, que sea corta. No rellenes.

EJEMPLOS DE CÓMO RESPONDER (sigue este estilo exacto):
- Si alguien dice "hola" → responde: "Hola [nombre si lo sabes]. ¿Qué necesitas?"
- Si alguien dice "gracias" → responde: "Para eso estoy."
- Si alguien dice "bien" → responde algo breve y natural, no una pregunta genérica
- Si alguien pide ayuda con un ejercicio → ve directo a resolverlo, sin introducción
- NUNCA respondas "¿En qué puedo ayudarte hoy?" — eso es lo que hace un robot, no un compañero"""

# ─────────────────────────────────────────────
# BASE DE DATOS
# ─────────────────────────────────────────────
def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS recordatorios (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, titulo TEXT, fecha_hora TEXT,
        mensaje TEXT, enviado INTEGER DEFAULT 0,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    c.execute("""CREATE TABLE IF NOT EXISTS contabilidad (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, tipo TEXT, concepto TEXT,
        monto REAL, fecha TEXT, pagado INTEGER DEFAULT 0,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    c.execute("""CREATE TABLE IF NOT EXISTS uso_mensual (
        user_id INTEGER, anio INTEGER, mes INTEGER,
        mensajes INTEGER DEFAULT 0, plan TEXT DEFAULT 'basico',
        PRIMARY KEY (user_id, anio, mes))""")
    c.execute("""CREATE TABLE IF NOT EXISTS horario_clases (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, materia TEXT, dia TEXT,
        hora_inicio TEXT, hora_fin TEXT, aula TEXT,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    c.execute("""CREATE TABLE IF NOT EXISTS examenes_pendientes (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, materia TEXT, fecha TEXT,
        hora TEXT, notificado INTEGER DEFAULT 0,
        analizado INTEGER DEFAULT 0,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    # ── TABLA CENTRAL DE USUARIOS ──
    c.execute("""CREATE TABLE IF NOT EXISTS usuarios (
        user_id INTEGER PRIMARY KEY,
        username TEXT, nombre TEXT,
        primera_vez TEXT DEFAULT CURRENT_TIMESTAMP,
        ultimo_mensaje TEXT DEFAULT CURRENT_TIMESTAMP)""")
    # ── SISTEMA DE SUSCRIPCIONES (pagos) ──
    c.execute("""CREATE TABLE IF NOT EXISTS suscripciones (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER UNIQUE,
        plan TEXT DEFAULT 'inactivo',
        fecha_inicio TEXT,
        fecha_vencimiento TEXT,
        aprobado_por INTEGER,
        notas TEXT,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    conn.commit()
    conn.close()


# ─────────────────────────────────────────────
# SISTEMA DE SUSCRIPCIONES
# ─────────────────────────────────────────────
ADMIN_ID = int(os.getenv("ADMIN_ID", "0"))  # Tu user_id de Telegram — setear en .env

def suscripcion_activa(user_id: int) -> bool:
    """Retorna True si el usuario tiene suscripción vigente."""
    # Admin siempre tiene acceso
    if user_id == ADMIN_ID:
        return True
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        SELECT plan, fecha_vencimiento FROM suscripciones
        WHERE user_id = ?
    """, (user_id,))
    row = c.fetchone()
    conn.close()
    if not row:
        return False
    plan, vencimiento = row
    if plan == "inactivo":
        return False
    if vencimiento:
        try:
            from datetime import timezone
            venc = datetime.fromisoformat(vencimiento)
            return datetime.now() < venc
        except Exception:
            return False
    return True  # sin fecha de vencimiento = vitalicio

def get_suscripcion(user_id: int) -> dict:
    """Retorna info de suscripción del usuario."""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("SELECT plan, fecha_inicio, fecha_vencimiento FROM suscripciones WHERE user_id=?", (user_id,))
    row = c.fetchone()
    conn.close()
    if not row:
        return {"plan": "inactivo", "activo": False}
    plan, inicio, vencimiento = row
    activo = suscripcion_activa(user_id)
    return {"plan": plan, "activo": activo, "inicio": inicio, "vencimiento": vencimiento}

def activar_suscripcion(user_id: int, plan: str, meses: int, admin_id: int, notas: str = "") -> str:
    """Activa o renueva suscripción de un usuario. Solo puede hacerlo el admin."""
    from datetime import timezone
    inicio = datetime.now()
    venc   = inicio + timedelta(days=30 * meses)
    conn   = sqlite3.connect(DB_PATH)
    c      = conn.cursor()
    c.execute("""
        INSERT INTO suscripciones (user_id, plan, fecha_inicio, fecha_vencimiento, aprobado_por, notas)
        VALUES (?, ?, ?, ?, ?, ?)
        ON CONFLICT(user_id) DO UPDATE SET
            plan=excluded.plan,
            fecha_inicio=excluded.fecha_inicio,
            fecha_vencimiento=excluded.fecha_vencimiento,
            aprobado_por=excluded.aprobado_por,
            notas=excluded.notas
    """, (user_id, plan, inicio.isoformat(), venc.isoformat(), admin_id, notas))
    conn.commit()
    conn.close()
    return venc.strftime("%d/%m/%Y")

def registrar_usuario(user_id: int, username: str):
    """Registra o actualiza usuario en la tabla central."""
    conn = sqlite3.connect(DB_PATH)
    c    = conn.cursor()
    c.execute("""
        INSERT INTO usuarios (user_id, username, ultimo_mensaje)
        VALUES (?, ?, CURRENT_TIMESTAMP)
        ON CONFLICT(user_id) DO UPDATE SET
            username=excluded.username,
            ultimo_mensaje=CURRENT_TIMESTAMP
    """, (user_id, username or ""))
    conn.commit()
    conn.close()

# ─────────────────────────────────────────────
# PERFIL DE USUARIO
# ─────────────────────────────────────────────
def cargar_perfil(user_id):
    path = f"{PERFILES_DIR}/{user_id}.json"
    if os.path.exists(path):
        with open(path) as f:
            return json.load(f)
    return {
        "nombre": None, "carrera": None, "ciclo": None,
        "universidad": None, "onboarding": 0, "plan": "basico",
        "tono": "normal"
    }

def guardar_perfil(user_id, perfil):
    with open(f"{PERFILES_DIR}/{user_id}.json", "w") as f:
        json.dump(perfil, f, ensure_ascii=False, indent=2)

# ─────────────────────────────────────────────
# HISTORIAL DE CONVERSACIÓN
# ─────────────────────────────────────────────
def cargar_historial(user_id):
    path = f"{HISTORIAL_DIR}/{user_id}.json"
    if os.path.exists(path):
        with open(path) as f:
            return json.load(f)
    return []

def guardar_historial(user_id, historial):
    with open(f"{HISTORIAL_DIR}/{user_id}.json", "w") as f:
        json.dump(historial, f, ensure_ascii=False)

def cargar_resumen(user_id):
    path = f"{HISTORIAL_DIR}/{user_id}_resumen.txt"
    if os.path.exists(path):
        with open(path) as f:
            return f.read()
    return ""

def guardar_resumen(user_id, resumen):
    with open(f"{HISTORIAL_DIR}/{user_id}_resumen.txt", "w") as f:
        f.write(resumen)

# ─────────────────────────────────────────────
# SYSTEM PROMPT DINÁMICO (según perfil)
# ─────────────────────────────────────────────
def construir_system_prompt(perfil):
    prompt = SYSTEM_PROMPT_BASE
    nombre = perfil.get("nombre")
    carrera = perfil.get("carrera")
    ciclo = perfil.get("ciclo")
    universidad = perfil.get("universidad")

    if nombre or carrera:
        prompt += f"\n\nCONTEXTO DEL ALUMNO:"
        if nombre:
            prompt += f"\n- Nombre: {nombre}"
        if carrera:
            prompt += f"\n- Carrera: {carrera}"
        if ciclo:
            prompt += f"\n- Ciclo: {ciclo}"
        if universidad:
            prompt += f"\n- Universidad: {universidad}"

    # Adaptar tono según ciclo
    try:
        ciclo_num = int(str(ciclo).replace("°", "").replace("to", "").replace("vo", "").strip()) if ciclo else 0
    except:
        ciclo_num = 0

    if ciclo_num in range(1, 4):
        prompt += "\n\nNIVEL: Alumno de ciclos iniciales. Explica cada concepto desde cero, usa analogías simples, más pasos intermedios."
    elif ciclo_num in range(4, 7):
        prompt += "\n\nNIVEL: Alumno de ciclo intermedio. Puedes asumir conocimiento de cálculo 1 y física básica. Nivel de detalle moderado."
    elif ciclo_num >= 7:
        prompt += "\n\nNIVEL: Alumno avanzado. Usa terminología técnica, asume bases sólidas, puedes ser más conciso en pasos obvios."

    # Contexto de carrera
    carrera_l = (carrera or "").lower()
    if "civil" in carrera_l:
        prompt += "\n\nCARRERA: Ingeniería Civil. Prioriza aplicaciones en estructuras, resistencia de materiales, hidráulica."
    elif "mecán" in carrera_l or "mecanic" in carrera_l:
        prompt += "\n\nCARRERA: Ingeniería Mecánica. Prioriza mecánica de fluidos, termodinámica, diseño mecánico."
    elif "electr" in carrera_l:
        prompt += "\n\nCARRERA: Ingeniería Eléctrica/Electrónica. Prioriza circuitos, señales, electromagnetismo."
    elif "sistem" in carrera_l or "inform" in carrera_l or "comput" in carrera_l:
        prompt += "\n\nCARRERA: Sistemas/Informática. Prioriza algoritmos, estructuras de datos, programación."
    elif "industri" in carrera_l:
        prompt += "\n\nCARRERA: Ingeniería Industrial. Prioriza estadística, investigación de operaciones, gestión."

    return prompt

# ─────────────────────────────────────────────
# CONTROL DE USO MENSUAL
# ─────────────────────────────────────────────
def verificar_limite(user_id, perfil):
    """
    Verifica si el usuario puede usar el bot.
    Primero chequea suscripción activa (tabla suscripciones).
    Si no tiene suscripción, bloquea el acceso.
    Admin siempre pasa.
    """
    # Admin siempre tiene acceso irrestricto
    if user_id == ADMIN_ID:
        return True, 99999

    # Verificar suscripción activa
    if not suscripcion_activa(user_id):
        return False, 0

    # Contar mensajes del mes (límite de uso para plan básico)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    ahora = datetime.now()
    c.execute("INSERT OR IGNORE INTO uso_mensual (user_id, anio, mes, mensajes, plan) VALUES (?,?,?,0,?)",
              (user_id, ahora.year, ahora.month, perfil.get("plan", "basico")))
    c.execute("SELECT mensajes, plan FROM uso_mensual WHERE user_id=? AND anio=? AND mes=?",
              (user_id, ahora.year, ahora.month))
    row = c.fetchone()
    conn.commit()
    conn.close()
    if not row:
        return True, 0
    mensajes, plan = row
    limite = PLAN_FULL_LIMITE if plan == "full" else PLAN_BASICO_LIMITE
    return mensajes < limite, mensajes

def incrementar_uso(user_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    ahora = datetime.now()
    c.execute("UPDATE uso_mensual SET mensajes=mensajes+1 WHERE user_id=? AND anio=? AND mes=?",
              (user_id, ahora.year, ahora.month))
    conn.commit()
    conn.close()

# ─────────────────────────────────────────────
# CONSTRUCCIÓN DEL PAYLOAD PARA LA API
# ─────────────────────────────────────────────
def construir_payload(user_id, mensaje_nuevo, perfil, sistema=None):
    historial = cargar_historial(user_id)
    resumen   = cargar_resumen(user_id)
    recientes = historial[-RECENT_MSGS:]
    sistema_final = sistema or construir_system_prompt(perfil)
    if resumen:
        sistema_final += f"\n\nRESUMEN DE CONVERSACIONES PREVIAS:\n{resumen}"
    msgs = [{"role": "system", "content": sistema_final}]
    msgs.extend(recientes)
    msgs.append({"role": "user", "content": mensaje_nuevo})
    return msgs, historial

# ─────────────────────────────────────────────
# GENERACIÓN DE RESUMEN AUTOMÁTICO
# ─────────────────────────────────────────────
async def generar_resumen(user_id, historial):
    if len(historial) < SUMMARY_EVERY:
        return
    try:
        texto = "\n".join([f"{m['role']}: {m['content'][:200]}" for m in historial[:-RECENT_MSGS]])
        loop = asyncio.get_event_loop()
        def llamar():
            return client.chat.completions.create(
                model=MODELO_RESUMEN,
                messages=[
                    {"role": "system", "content": "Haz un resumen breve (máx 200 palabras) de esta conversación académica. Destaca: temas vistos, nivel del alumno, pendientes importantes."},
                    {"role": "user", "content": texto}
                ],
                max_tokens=300
            )
        resp = await loop.run_in_executor(None, llamar)
        guardar_resumen(user_id, resp.choices[0].message.content)
        guardar_historial(user_id, historial[-RECENT_MSGS:])
    except Exception as e:
        logger.error(f"Error generando resumen: {e}")

# ─────────────────────────────────────────────
# DETECCIÓN INTELIGENTE DE INTENCIÓN
# ─────────────────────────────────────────────
KEYWORDS_EJERCICIO_DIFICIL = [
    "integral", "derivada", "límite", "transformada", "laplace", "fourier",
    "diferencial", "ecuación diferencial", "cálculo", "mecánica de fluidos",
    "termodinámica", "resistencia de materiales", "momento de inercia",
    "estática", "dinámica", "cinemática", "cinematica", "vectores",
    "electromagnet", "circuito", "señales", "álgebra lineal", "matriz",
    "determinante", "eigenvalor", "autovalor", "series de taylor",
    "series de fourier", "variable compleja", "probabilidad avanzada",
    "estadística inferencial", "regresión", "mecánica cuántica",
    "termodinamica", "fluidos", "hidráulica", "hidraulica",
    "torsión", "flexión", "pandeo", "esfuerzo", "deformacion",
    "deformación", "mohr", "bernoulli", "navier", "reynolds"
]

KEYWORDS_EJERCICIO_SIMPLE = [
    "resolver", "calcular", "encontrar", "hallar", "determinar",
    "ejercicio", "problema", "tarea", "practica", "práctica",
    "suma", "resta", "multiplica", "divide", "porcentaje",
    "velocidad", "aceleración", "fuerza", "energia", "energía",
    "trabajo", "potencia", "presión", "presion", "densidad",
    "temperatura", "volumen", "masa", "peso", "newton"
]

KEYWORDS_EXCEL = [
    "excel", "xlsx", "hoja de cálculo", "hoja de calculo",
    "tabla de datos", "tabla de notas", "registro", "seguimiento",
    "presupuesto", "gastos en tabla", "datos del laboratorio",
    "tabla de gastos", "planilla"
]

KEYWORDS_PPTX = [
    "presentación", "presentacion", "diapositiva", "diapositivas",
    "slides", "ppt", "pptx", "exposición", "exposicion", "sustentación"
]

KEYWORDS_DOCX = [
    "informe", "reporte word", "documento word", "docx",
    "carta", "redactar informe", "informe de laboratorio",
    "resumen en documento", "formato word"
]

KEYWORDS_RECORDATORIO = [
    "recuérdame", "recuerdame", "avísame", "avisame",
    "recordar", "recordatorio", "no olvides", "a las", "el lunes",
    "el martes", "el miércoles", "el jueves", "el viernes",
    "mañana a", "pasado mañana", "el domingo", "el sábado",
    "examen el", "clase el", "reunión el", "pagar el"
]

KEYWORDS_CONTABILIDAD = [
    "gasté", "gaste", "compré", "compre", "pagué", "pague",
    "me depositaron", "cobré", "cobre", "ingresé", "ingrese",
    "gasto de", "ingreso de", "mi saldo", "mi balance",
    "cuánto llevo gastado", "mis gastos", "mis ingresos",
    "cuanto llevo", "cuánto tengo"
]

KEYWORDS_HORARIO = [
    "mi horario", "ver horario", "clases hoy", "que tengo hoy",
    "qué tengo hoy", "resumen del dia", "resumen del día",
    "mis clases", "horario de hoy", "agenda de hoy",
    "tengo clases", "a qué hora", "a que hora tengo"
]

KEYWORDS_EXAMEN = [
    "tengo examen", "examen de", "registrar examen", "examen el",
    "hay examen", "fecha de examen", "mi examen"
]

KEYWORDS_ENTREGA = [
    "entrega de", "entregar", "presentar", "fecha de entrega",
    "tengo que entregar", "plazo de"
]

def detectar_intencion(texto):
    t = texto.lower()

    # Recordatorio (prioridad alta — tiene palabras muy específicas)
    if any(k in t for k in KEYWORDS_RECORDATORIO):
        return "recordatorio"

    # Horario / agenda académica
    if any(k in t for k in KEYWORDS_HORARIO):
        return "horario"

    # Examen
    if any(k in t for k in KEYWORDS_EXAMEN):
        return "examen"

    # Entrega
    if any(k in t for k in KEYWORDS_ENTREGA):
        return "entrega"

    # Contabilidad
    if any(k in t for k in KEYWORDS_CONTABILIDAD):
        return "contabilidad"

    # Archivos
    if any(k in t for k in KEYWORDS_EXCEL):
        return "excel"
    if any(k in t for k in KEYWORDS_PPTX):
        return "presentacion"
    if any(k in t for k in KEYWORDS_DOCX):
        return "documento"

    # Ejercicios
    if any(k in t for k in KEYWORDS_EJERCICIO_DIFICIL):
        return "ejercicio_dificil"
    if any(k in t for k in KEYWORDS_EJERCICIO_SIMPLE):
        return "ejercicio_simple"

    return "chat"

def elegir_modelo(intencion):
    mapa = {
        "ejercicio_dificil":  MODELO_EJERCICIOS,
        "ejercicio_simple":   MODELO_EJERCICIOS_RAPIDO,
        "excel":              MODELO_CHAT,
        "presentacion":       MODELO_CHAT,
        "documento":          MODELO_CHAT,
        "chat":               MODELO_CHAT,
        "recordatorio":       MODELO_RAPIDO,
        "contabilidad":       MODELO_RAPIDO,
        "horario":            MODELO_RAPIDO,
        "examen":             MODELO_RAPIDO,
        "entrega":            MODELO_RAPIDO,
    }
    return mapa.get(intencion, MODELO_CHAT)


# ─────────────────────────────────────────────
# PROCESADOR DE LATEX → TEXTO LEGIBLE
# ─────────────────────────────────────────────
LATEX_MAP = {
    r'\alpha': 'α', r'\beta': 'β', r'\gamma': 'γ', r'\delta': 'δ',
    r'\epsilon': 'ε', r'\theta': 'θ', r'\lambda': 'λ', r'\mu': 'μ',
    r'\nu': 'ν', r'\pi': 'π', r'\rho': 'ρ', r'\sigma': 'σ',
    r'\tau': 'τ', r'\phi': 'φ', r'\psi': 'ψ', r'\omega': 'ω',
    r'\Omega': 'Ω', r'\Delta': 'Δ', r'\Sigma': 'Σ', r'\Pi': 'Π',
    r'\infty': '∞', r'\partial': '∂', r'\nabla': '∇',
    r'\int': '∫', r'\sum': '∑', r'\prod': '∏',
    r'\leq': '≤', r'\geq': '≥', r'\neq': '≠', r'\approx': '≈',
    r'\pm': '±', r'\times': '×', r'\div': '÷', r'\cdot': '·',
    r'\sqrt': '√', r'\degree': '°',
    r'\rightarrow': '→', r'\leftarrow': '←', r'\Rightarrow': '⇒',
    r'\ldots': '...', r'\cdots': '···',
}

def limpiar_latex(texto):
    """Convierte LaTeX crudo a texto legible — limpieza agresiva."""
    if not texto:
        return texto

    # Eliminar bloques de código LaTeX de documento completo
    texto = re.sub(r'```(?:tex|latex).*?```', '', texto, flags=re.DOTALL)
    texto = re.sub(r'\\documentclass.*?\\end\{document\}', '', texto, flags=re.DOTALL)
    texto = re.sub(r'\\documentclass[^\n]*\n', '', texto)
    texto = re.sub(r'\\usepackage[^\n]*\n', '', texto)
    texto = re.sub(r'\\geometry[^\n]*\n', '', texto)
    texto = re.sub(r'\\begin\{document\}[^\n]*\n?', '', texto)
    texto = re.sub(r'\\end\{document\}[^\n]*\n?', '', texto)
    texto = re.sub(r'\\maketitle[^\n]*\n?', '', texto)
    texto = re.sub(r'\\author\{[^}]*\}', '', texto)
    texto = re.sub(r'\\date\{[^}]*\}', '', texto)
    texto = re.sub(r'\\title\{([^}]*)\}', r'\1', texto)
    texto = re.sub(r'\\section\*?\{([^}]*)\}', r'\n### \1\n', texto)
    texto = re.sub(r'\\subsection\*?\{([^}]*)\}', r'\n**\1**\n', texto)

    # Letras griegas y símbolos (LATEX_MAP)
    for latex, unicode_char in LATEX_MAP.items():
        texto = texto.replace(latex, unicode_char)

    # Fórmulas: \frac{a}{b} → (a/b)
    texto = re.sub(r'\\frac\{([^}]+)\}\{([^}]+)\}', r'(\1/\2)', texto)
    texto = re.sub(r'\\sqrt\{([^}]+)\}', r'√(\1)', texto)
    texto = re.sub(r'\\sqrt\b', '√', texto)

    # Superíndices
    sup_map = {'0':'⁰','1':'¹','2':'²','3':'³','4':'⁴','5':'⁵',
               '6':'⁶','7':'⁷','8':'⁸','9':'⁹','n':'ⁿ'}
    def repl_sup(m):
        return ''.join(sup_map.get(ch, ch) for ch in m.group(1))
    texto = re.sub(r'\^\{([^}]+)\}', repl_sup, texto)
    texto = re.sub(r'\^(\d)', lambda m: sup_map.get(m.group(1), m.group(1)), texto)

    # Subíndices
    texto = re.sub(r'_\{([^}]+)\}', r'_\1', texto)

    # Delimitadores matemáticos
    texto = re.sub(r'\\\[(.+?)\\\]', r'\1', texto, flags=re.DOTALL)
    texto = re.sub(r'\\\((.+?)\\\)', r'\1', texto, flags=re.DOTALL)
    texto = re.sub(r'\$\$(.+?)\$\$', r'\1', texto, flags=re.DOTALL)
    texto = re.sub(r'\$(.+?)\$', r'\1', texto)

    # Comandos de formato
    texto = re.sub(r'\\text\{([^}]+)\}', r'\1', texto)
    texto = re.sub(r'\\(bf|it|rm|mathbf|mathit|mathrm|mathbb|textbf|textit)\{([^}]+)\}', r'\2', texto)
    texto = re.sub(r'\\(left|right|big|Big|bigg|Bigg)[()\[\]|{}]?', '', texto)

    # Eliminar comandos LaTeX sueltos restantes
    texto = re.sub(r'\\[a-zA-Z]+\*?(?:\{[^}]*\})*', '', texto)

    # Limpiar llaves y símbolos sobrantes
    texto = texto.replace('{', '').replace('}', '')
    texto = re.sub(r' {2,}', ' ', texto)
    texto = re.sub(r'\n{3,}', '\n\n', texto)

    return texto.strip()


# ─────────────────────────────────────────────
# RENDERIZADOR DE FÓRMULAS MATEMÁTICAS (matplotlib)
# ─────────────────────────────────────────────
def renderizar_formula(formula_latex, ancho_px=500, alto_px=80, fontsize=16):
    """
    Renderiza una fórmula LaTeX como imagen PNG usando matplotlib.mathtext.
    Retorna bytes de la imagen o None si falla.
    """
    try:
        fig, ax = plt.subplots(figsize=(ancho_px/100, alto_px/100), dpi=100)
        ax.axis('off')
        # Si no tiene $ ya, los agregamos
        formula_display = formula_latex if formula_latex.startswith('$') else f'${formula_latex}$'
        ax.text(0.5, 0.5, formula_display,
                ha='center', va='center',
                fontsize=fontsize,
                transform=ax.transAxes,
                color='#1a1a2e')
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor='white', edgecolor='none', pad_inches=0.1)
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        logger.warning(f"Error renderizando fórmula '{formula_latex}': {e}")
        plt.close('all')
        return None


# ─────────────────────────────────────────────
# GENERADOR DE PDF PROFESIONAL CON MATEMÁTICAS
# ─────────────────────────────────────────────
# ─────────────────────────────────────────────
# PDF — usa agente_academico (arquitectura JSON)
# ─────────────────────────────────────────────
def _crear_pdf_limpio(contenido, user_id, titulo="Solución", perfil=None):
    """Wrapper: pasa contenido al agente académico para generar PDF."""
    from modulos.agente_academico import crear_pdf_desde_json, SYSTEM_JSON
    import json as _json, re as _re
    # Intentar parsear si ya es JSON
    try:
        txt = contenido.strip()
        txt = _re.sub(r'^```json\s*|^```\s*|```$', '', txt, flags=_re.MULTILINE).strip()
        datos = _json.loads(txt)
        return crear_pdf_desde_json(datos, user_id, titulo, perfil)
    except Exception:
        # Si no es JSON, crear estructura básica con el texto
        datos = {"ejercicios": [{"titulo": titulo, "datos": [], "pasos": [{"num": 1, "titulo": "Solución", "calculo": contenido[:2000]}], "resultado": ""}]}
        return crear_pdf_desde_json(datos, user_id, titulo, perfil)


# ─────────────────────────────────────────────
def crear_excel(contenido_json, user_id, titulo="Datos"):
    """
    Genera un archivo Excel profesional con:
    - Headers formateados (azul oscuro, texto blanco)
    - Fórmulas Excel reales (no valores hardcodeados)
    - Filas alternadas para lectura
    - Columnas con ancho automático
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{ARCHIVOS_DIR}/datos_{user_id}_{timestamp}.xlsx"

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = titulo[:30]

    # Colores profesionales
    COLOR_HEADER  = "1a1a2e"
    COLOR_ALT_ROW = "eef2ff"
    COLOR_ACCENT  = "0f3460"
    COLOR_TOTAL   = "e94560"

    borde_fino = Border(
        left=Side(style='thin', color="cccccc"),
        right=Side(style='thin', color="cccccc"),
        top=Side(style='thin', color="cccccc"),
        bottom=Side(style='thin', color="cccccc")
    )

    # contenido_json puede ser dict con 'headers', 'rows', 'totales'
    headers = contenido_json.get("headers", [])
    rows    = contenido_json.get("rows", [])
    totales = contenido_json.get("totales", False)

    # Título del documento
    ws.merge_cells(f"A1:{get_column_letter(max(len(headers), 1))}1")
    celda_titulo = ws["A1"]
    celda_titulo.value = titulo
    celda_titulo.font = Font(name="Calibri", bold=True, size=14, color="FFFFFF")
    celda_titulo.fill = PatternFill("solid", fgColor=COLOR_ACCENT)
    celda_titulo.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28

    # Headers
    for col, header in enumerate(headers, 1):
        c = ws.cell(row=2, column=col, value=header)
        c.font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor=COLOR_HEADER)
        c.alignment = Alignment(horizontal="center", vertical="center")
        c.border = borde_fino
    ws.row_dimensions[2].height = 22

    # Datos
    for fila_idx, fila in enumerate(rows, 3):
        bg = COLOR_ALT_ROW if fila_idx % 2 == 1 else "FFFFFF"
        for col, valor in enumerate(fila, 1):
            c = ws.cell(row=fila_idx, column=col, value=valor)
            c.font = Font(name="Calibri", size=10)
            c.fill = PatternFill("solid", fgColor=bg)
            c.alignment = Alignment(vertical="center", wrap_text=True)
            c.border = borde_fino

    # Fila de totales (usa fórmulas Excel reales)
    if totales and rows:
        fila_total = len(rows) + 3
        ws.cell(row=fila_total, column=1, value="TOTAL").font = Font(
            name="Calibri", bold=True, size=11, color="FFFFFF")
        ws.cell(row=fila_total, column=1).fill = PatternFill("solid", fgColor=COLOR_TOTAL)

        for col in range(2, len(headers) + 1):
            letra = get_column_letter(col)
            c = ws.cell(row=fila_total, column=col,
                        value=f"=SUM({letra}3:{letra}{fila_total-1})")
            c.font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
            c.fill = PatternFill("solid", fgColor=COLOR_TOTAL)
            c.number_format = '#,##0.00'
            c.border = borde_fino

    # Ancho automático de columnas
    for col in range(1, len(headers) + 1):
        max_len = 0
        letra = get_column_letter(col)
        for row in ws.iter_rows(min_col=col, max_col=col):
            for cell in row:
                try:
                    max_len = max(max_len, len(str(cell.value or "")))
                except:
                    pass
        ws.column_dimensions[letra].width = min(max(max_len + 4, 12), 40)

    # Congelar primera fila de datos
    ws.freeze_panes = "A3"

    wb.save(path)
    return path


# ─────────────────────────────────────────────
# GENERADOR DE PRESENTACIÓN PPTX PROFESIONAL
# ─────────────────────────────────────────────
def crear_presentacion(slides_data, user_id, titulo_pres="Presentación"):
    """
    Genera PPTX con diseño profesional.
    slides_data: lista de dicts con 'titulo', 'contenido', 'tipo'
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{ARCHIVOS_DIR}/presentacion_{user_id}_{timestamp}.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Paleta — Midnight Executive
    C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
    C_MID    = RGBColor(0x16, 0x21, 0x3e)
    C_ACCENT = RGBColor(0x0f, 0x34, 0x60)
    C_GOLD   = RGBColor(0xe9, 0x45, 0x60)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_LIGHT  = RGBColor(0xCA, 0xDC, 0xFC)

    def add_bg(slide, color):
        from pptx.util import Emu
        bg = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        # Mover al fondo
        sp_tree = slide.shapes._spTree
        sp_tree.remove(bg._element)
        sp_tree.insert(2, bg._element)

    def add_text(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        from pptx.util import Pt
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or C_WHITE
        run.font.name = "Calibri"
        return txBox

    def add_rect(slide, left, top, width, height, color, radius=False):
        shape_type = 5 if radius else 1  # 5=rounded rect, 1=rect
        s = slide.shapes.add_shape(shape_type, left, top, width, height)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    layout = prs.slide_layouts[6]  # En blanco

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(layout)
        tipo = slide_info.get("tipo", "contenido")
        titulo_sl = slide_info.get("titulo", "")
        contenido = slide_info.get("contenido", "")

        W = prs.slide_width
        H = prs.slide_height

        if tipo == "portada":
            # Fondo oscuro completo
            add_bg(slide, C_DARK)
            # Barra de acento lateral
            add_rect(slide, 0, 0, Inches(0.5), H, C_GOLD)
            # Título grande
            add_text(slide, titulo_sl,
                     Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                     font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
            # Subtítulo
            if contenido:
                add_text(slide, contenido,
                         Inches(1), Inches(3.9), Inches(10), Inches(1),
                         font_size=20, color=C_LIGHT, align=PP_ALIGN.LEFT)
            # Fecha
            add_text(slide, datetime.now().strftime("%B %Y"),
                     Inches(1), Inches(5.2), Inches(5), Inches(0.5),
                     font_size=14, color=C_LIGHT)
            # Número de slide
            add_text(slide, "JARVIS 3.0",
                     Inches(10), Inches(6.8), Inches(3), Inches(0.5),
                     font_size=10, color=C_LIGHT, align=PP_ALIGN.RIGHT)

        elif tipo == "cierre":
            add_bg(slide, C_DARK)
            add_rect(slide, 0, 0, Inches(0.5), H, C_GOLD)
            add_text(slide, titulo_sl,
                     Inches(1), Inches(2.8), Inches(11), Inches(1.2),
                     font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            if contenido:
                add_text(slide, contenido,
                         Inches(1), Inches(4.2), Inches(11), Inches(1),
                         font_size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)

        else:  # contenido estándar
            add_bg(slide, RGBColor(0xF8, 0xF9, 0xFF))  # Fondo muy claro
            # Barra superior oscura
            add_rect(slide, 0, 0, W, Inches(1.1), C_DARK)
            # Título en barra
            add_text(slide, titulo_sl,
                     Inches(0.4), Inches(0.15), Inches(11), Inches(0.8),
                     font_size=26, bold=True, color=C_WHITE)
            # Número de slide (barra lateral derecha)
            add_rect(slide, W - Inches(0.6), 0, Inches(0.6), Inches(1.1), C_GOLD)
            add_text(slide, str(idx + 1),
                     W - Inches(0.6), Inches(0.2), Inches(0.6), Inches(0.7),
                     font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

            # Contenido como puntos
            puntos = [p.strip() for p in contenido.split('\n') if p.strip()]
            top_offset = Inches(1.3)
            for i, punto in enumerate(puntos[:7]):
                # Bullet circle
                add_rect(slide, Inches(0.4), top_offset + Inches(0.15),
                         Inches(0.22), Inches(0.22), C_GOLD)
                # Texto
                add_text(slide, punto,
                         Inches(0.75), top_offset, Inches(12.1), Inches(0.7),
                         font_size=16, color=C_DARK)
                top_offset += Inches(0.78)

            # Línea decorativa inferior
            add_rect(slide, Inches(0.4), H - Inches(0.4), Inches(12.5), Inches(0.06), C_ACCENT)

    prs.save(path)
    return path


# ─────────────────────────────────────────────
# GENERADOR DE DOCUMENTO WORD PROFESIONAL
# ─────────────────────────────────────────────
def crear_documento_word(contenido, user_id, titulo="Documento"):
    """
    Genera un documento .docx profesional con python-docx.
    """
    from docx import Document
    from docx.shared import Pt, Cm, RGBColor as DocxRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{ARCHIVOS_DIR}/documento_{user_id}_{timestamp}.docx"

    doc = Document()

    # Márgenes
    section = doc.sections[0]
    section.top_margin    = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin   = Cm(3)
    section.right_margin  = Cm(2.5)

    # Título principal
    titulo_p = doc.add_heading(titulo, level=0)
    titulo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = titulo_p.runs[0]
    run.font.color.rgb = DocxRGB(0x1a, 0x1a, 0x2e)
    run.font.size = Pt(22)

    # Subtítulo con fecha
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub.add_run(datetime.now().strftime("%d de %B de %Y"))
    sub_run.font.color.rgb = DocxRGB(0x66, 0x66, 0x88)
    sub_run.font.size = Pt(11)

    # Línea separadora
    doc.add_paragraph("─" * 60)

    # Procesar contenido
    for linea in contenido.split('\n'):
        linea_limpia = limpiar_latex(linea.strip())
        if not linea_limpia:
            doc.add_paragraph()
            continue

        if linea_limpia.startswith('## '):
            h = doc.add_heading(linea_limpia[3:], level=2)
            h.runs[0].font.color.rgb = DocxRGB(0x0f, 0x34, 0x60)
        elif linea_limpia.startswith('# '):
            h = doc.add_heading(linea_limpia[2:], level=1)
            h.runs[0].font.color.rgb = DocxRGB(0x1a, 0x1a, 0x2e)
        elif linea_limpia.startswith('- ') or linea_limpia.startswith('• '):
            p = doc.add_paragraph(style='List Bullet')
            p.add_run(linea_limpia[2:]).font.size = Pt(11)
        elif re.match(r'^\d+\.\s', linea_limpia):
            p = doc.add_paragraph(style='List Number')
            p.add_run(re.sub(r'^\d+\.\s', '', linea_limpia)).font.size = Pt(11)
        else:
            p = doc.add_paragraph(linea_limpia)
            p.runs[0].font.size = Pt(11) if p.runs else None

    # Pie de página
    footer = section.footer
    footer_p = footer.paragraphs[0]
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_p.add_run(
        f"Generado por JARVIS 3.0 — {datetime.now().strftime('%d/%m/%Y')}")
    footer_run.font.size = Pt(8)
    footer_run.font.color.rgb = DocxRGB(0x99, 0x99, 0x99)

    doc.save(path)
    return path


# ─────────────────────────────────────────────
# PARSER INTELIGENTE DE ARCHIVOS (IA genera la estructura)
# ─────────────────────────────────────────────
async def generar_estructura_archivo(tipo, solicitud, perfil):
    """
    Usa la IA para generar la estructura de datos del archivo solicitado.
    Retorna un dict con la estructura lista para usar.
    """
    prompts = {
        "excel": f"""El usuario solicita: "{solicitud}"
Genera la estructura de datos para un Excel profesional.
Responde SOLO con JSON válido, sin texto extra, sin bloques de código:
{{
  "titulo": "Nombre del Excel",
  "headers": ["Col1", "Col2", ...],
  "rows": [["val1", "val2", ...], ...],
  "totales": true/false
}}
Incluye datos de ejemplo realistas si el usuario no proporcionó datos.
Los valores numéricos deben ser números (no strings).""",

        "presentacion": f"""El usuario solicita: "{solicitud}"
Genera la estructura para una presentación profesional.
Responde SOLO con JSON válido, sin texto extra, sin bloques de código:
{{
  "titulo": "Título de la presentación",
  "slides": [
    {{"tipo": "portada", "titulo": "Título", "contenido": "subtítulo"}},
    {{"tipo": "contenido", "titulo": "Slide 2", "contenido": "punto 1\\npunto 2\\npunto 3"}},
    {{"tipo": "cierre", "titulo": "Conclusiones", "contenido": "mensaje final"}}
  ]
}}
Mínimo 5 slides, máximo 12. Incluye portada y cierre siempre.""",

        "documento": f"""El usuario solicita: "{solicitud}"
Genera el contenido para un documento Word profesional.
Responde SOLO con JSON válido, sin texto extra, sin bloques de código:
{{
  "titulo": "Título del documento",
  "contenido": "# Sección 1\\nTexto del párrafo...\\n## Subsección\\n- Punto 1\\n- Punto 2"
}}
Usa # para H1, ## para H2, - para bullets, números para listas."""
    }

    loop = asyncio.get_event_loop()
    prompt_usado = prompts.get(tipo, prompts["documento"])

    def llamar():
        return client.chat.completions.create(
            model=MODELO_CHAT,
            messages=[
                {"role": "system", "content": "Eres un generador de estructuras de datos. Responde ÚNICAMENTE con JSON válido. Cero texto adicional."},
                {"role": "user", "content": prompt_usado}
            ],
            max_tokens=2000
        )

    try:
        resp = await loop.run_in_executor(None, llamar)
        texto = resp.choices[0].message.content.strip()
        # Limpiar posibles bloques de código
        texto = re.sub(r'^```(?:json)?\s*', '', texto)
        texto = re.sub(r'\s*```$', '', texto)
        return json.loads(texto)
    except Exception as e:
        logger.error(f"Error generando estructura {tipo}: {e}")
        return None


# ─────────────────────────────────────────────
# RECORDATORIOS — PARSER DE FECHA NATURAL
# ─────────────────────────────────────────────
DIAS_SEMANA = {
    "lunes": 0, "martes": 1, "miércoles": 2, "miercoles": 2,
    "jueves": 3, "viernes": 4, "sábado": 5, "sabado": 5, "domingo": 6
}

def parsear_fecha_natural(texto):
    """
    Extrae fecha y hora de un texto en lenguaje natural.
    Retorna datetime o None.
    """
    texto = texto.lower()
    ahora = datetime.now()

    # "mañana a las HH:MM" o "mañana a las H pm"
    m = re.search(r'ma[ñn]ana\s+a\s+las?\s+(\d{1,2})(?::(\d{2}))?\s*(am|pm)?', texto)
    if m:
        hora, minuto, ampm = m.group(1), m.group(2) or "0", m.group(3) or ""
        hora = int(hora)
        if ampm == "pm" and hora < 12: hora += 12
        elif ampm == "am" and hora == 12: hora = 0
        return (ahora + timedelta(days=1)).replace(hour=hora, minute=int(minuto), second=0, microsecond=0)

    # "el [día] a las HH:MM"
    for dia, num_dia in DIAS_SEMANA.items():
        if dia in texto:
            m = re.search(r'a\s+las?\s+(\d{1,2})(?::(\d{2}))?\s*(am|pm)?', texto)
            hora = int(m.group(1)) if m else 8
            minuto = int(m.group(2)) if m and m.group(2) else 0
            ampm = m.group(3) if m else ""
            if ampm == "pm" and hora < 12: hora += 12
            elif ampm == "am" and hora == 12: hora = 0
            dias_diff = (num_dia - ahora.weekday()) % 7
            if dias_diff == 0: dias_diff = 7
            return (ahora + timedelta(days=dias_diff)).replace(
                hour=hora, minute=minuto, second=0, microsecond=0)

    # "en X horas" / "en X minutos"
    m = re.search(r'en\s+(\d+)\s+(hora|horas|minuto|minutos)', texto)
    if m:
        cantidad, unidad = int(m.group(1)), m.group(2)
        if "hora" in unidad:
            return ahora + timedelta(hours=cantidad)
        else:
            return ahora + timedelta(minutes=cantidad)

    # "a las HH:MM" (hoy)
    m = re.search(r'a\s+las?\s+(\d{1,2})(?::(\d{2}))?\s*(am|pm)?', texto)
    if m:
        hora = int(m.group(1))
        minuto = int(m.group(2) or "0")
        ampm = m.group(3) or ""
        if ampm == "pm" and hora < 12: hora += 12
        elif ampm == "am" and hora == 12: hora = 0
        dt = ahora.replace(hour=hora, minute=minuto, second=0, microsecond=0)
        if dt <= ahora: dt += timedelta(days=1)
        return dt

    return None


def extraer_titulo_recordatorio(texto):
    """Extrae el título del evento del texto de recordatorio."""
    texto = re.sub(r'recuérdame|recuerdame|avísame|avisame|recordar', '', texto, flags=re.IGNORECASE)
    texto = re.sub(r'a\s+las?\s+\d{1,2}(?::\d{2})?\s*(?:am|pm)?', '', texto, flags=re.IGNORECASE)
    texto = re.sub(r'el\s+(?:lunes|martes|miércoles|miercoles|jueves|viernes|sábado|sabado|domingo)', '', texto, flags=re.IGNORECASE)
    texto = re.sub(r'ma[ñn]ana|pasado\s+ma[ñn]ana|en\s+\d+\s+(?:hora|minuto)s?', '', texto, flags=re.IGNORECASE)
    texto = re.sub(r'[,\.\-]+', ' ', texto)
    return texto.strip().capitalize() or "Recordatorio"


# ─────────────────────────────────────────────
# VISIÓN — PIPELINE DE 2 PASOS
# ─────────────────────────────────────────────
async def analizar_imagen_completo(image_bytes, caption, user_id, perfil):
    """
    Pipeline de 3 pasos:
    1. llama-4-scout: lee la imagen y transcribe TODOS los ejercicios
    2. gpt-oss-120b: resuelve TODOS los ejercicios uno por uno
    3. llama-3.1-8b: convierte LaTeX a texto limpio
    """
    loop = asyncio.get_event_loop()
    img_b64 = base64.b64encode(image_bytes).decode("utf-8")

    # ── PASO 1: Scout lee y transcribe todos los ejercicios ──
    def extraer_contexto():
        return client.chat.completions.create(
            model=MODELO_VISION,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}
                    },
                    {
                        "type": "text",
                        "text": (
                            "Transcribe EXACTAMENTE todo el contenido de esta imagen.\n"
                            "Si hay ejercicios numerados, transcríbelos TODOS con su número.\n"
                            "Incluye: enunciados, datos, fórmulas, opciones de respuesta, todo.\n"
                            "No omitas ningún ejercicio. No resuelvas nada, solo transcribe."
                        )
                    }
                ]
            }],
            max_tokens=2000
        )

    try:
        ctx_resp = await loop.run_in_executor(None, extraer_contexto)
        contexto = ctx_resp.choices[0].message.content
        logger.info(f"Scout transcribió: {len(contexto)} chars")
    except Exception as e:
        logger.error(f"Error extracción visual: {e}")
        contexto = caption or "Imagen enviada por el alumno"

    # ── PASO 2: 120b resuelve TODOS los ejercicios ──
    _cap = re.sub(
        r'\b(en\s+)?pdf\b|de\s+manera\s+profesional|formato\s+pdf|como\s+pdf|resuélvelo|resuelvelo',
        '', caption, flags=re.IGNORECASE
    ).strip(" ,.-")
    caption_extra = f"\nNota del alumno: {_cap}" if _cap else ""

    def resolver():
        return client.chat.completions.create(
            model=MODELO_EJERCICIOS,  # gpt-oss-120b
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Eres un tutor experto en ciencias e ingeniería. "
                        "Resuelves ejercicios de forma completa, paso a paso, con todos los cálculos."
                    )
                },
                {
                    "role": "user",
                    "content": (
                        f"Aquí están los ejercicios transcritos de una imagen:{caption_extra}\n\n"
                        f"{contexto}\n\n"
                        f"TAREA: Resuelve CADA ejercicio que aparece arriba, uno por uno, en orden.\n"
                        f"No omitas ninguno. Si aparecen 3 ejercicios, resuelve los 3.\n"
                        f"Para cada ejercicio:\n"
                        f"- Escribe el número del ejercicio\n"
                        f"- Escribe los datos conocidos\n"
                        f"- Muestra el procedimiento completo con todos los cálculos\n"
                        f"- Escribe la respuesta final con unidades\n"
                        f"Usa texto plano para las fórmulas, no LaTeX."
                    )
                }
            ],
            max_tokens=5000
        )

    try:
        sol_resp = await loop.run_in_executor(None, resolver)
        solucion_raw = sol_resp.choices[0].message.content
        logger.info(f"120b resolvió: {len(solucion_raw)} chars")
    except Exception as e:
        logger.error(f"Error resolviendo con 120b: {e}")
        return contexto, "No pude resolver los ejercicios. Intenta describírmelos en texto.", MODELO_EJERCICIOS

    # ── PASO 3: llama-3.1-8b limpia todo el LaTeX restante ──
    def limpiar_con_ia():
        return client.chat.completions.create(
            model=MODELO_RAPIDO,  # llama-3.1-8b-instant
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Eres un limpiador de texto matemático. Tu única tarea:\n"
                        "Recibir una solución y reescribirla sin ningún LaTeX.\n"
                        "Conversiones obligatorias:\n"
                        "- \\frac{a}{b} → a/b\n"
                        "- \\sqrt{x} → √x\n"
                        "- x^{2} → x²\n"
                        "- \\alpha → α, \\beta → β, \\Delta → Δ, \\mu → μ\n"
                        "- \\times → ×, \\cdot → ·\n"
                        "- $$...$$ y $...$ → solo el contenido sin signos $\n"
                        "- \\[...\\] y \\(...\\) → solo el contenido\n"
                        "- Eliminar: \\documentclass, \\begin{}, \\end{}, \\usepackage, \\section, \\maketitle\n"
                        "Mantén todos los números, procedimientos y respuestas intactos.\n"
                        "Formato final: cada ejercicio separado por línea en blanco, numerado.\n"
                        "No añadas ni quites contenido matemático, solo limpia el formato."
                    )
                },
                {
                    "role": "user",
                    "content": f"Limpia esta solución:\n\n{solucion_raw}"
                }
            ],
            max_tokens=5000
        )

    try:
        clean_resp = await loop.run_in_executor(None, limpiar_con_ia)
        solucion = clean_resp.choices[0].message.content
        logger.info(f"8b limpió: {len(solucion)} chars")
    except Exception as e:
        logger.error(f"Error limpiando con 8b: {e}")
        solucion = limpiar_latex(solucion_raw)  # fallback

    return contexto, solucion, MODELO_EJERCICIOS


# ─────────────────────────────────────────────
# CONTABILIDAD — PARSER INTELIGENTE
# ─────────────────────────────────────────────
def parsear_movimiento(texto):
    """Extrae tipo, monto y concepto de un texto de contabilidad."""
    texto_l = texto.lower()
    tipo = "gasto"
    if any(k in texto_l for k in ["depositaron", "cobré", "cobre", "me pagaron",
                                    "ingresé", "ingrese", "recibí", "recibi", "ingreso"]):
        tipo = "ingreso"

    # Extraer monto
    m = re.search(r'(\d+(?:\.\d{1,2})?)\s*(?:soles?|sol|s/|pen)?', texto_l)
    monto = float(m.group(1)) if m else 0.0

    # Extraer concepto (eliminar palabras de trigger y monto)
    concepto = re.sub(r'gasté|gaste|compré|compre|pagué|pague|me depositaron|cobré|ingresé|ingrese', '', texto, flags=re.IGNORECASE)
    concepto = re.sub(r'\d+(?:\.\d{1,2})?\s*(?:soles?|sol|s/)?', '', concepto, flags=re.IGNORECASE)
    concepto = re.sub(r'[,\.\-]+', ' ', concepto).strip().capitalize()
    if not concepto:
        concepto = "Sin descripción"

    return tipo, monto, concepto


# ─────────────────────────────────────────────
# VOZ — TRANSCRIPCIÓN
# ─────────────────────────────────────────────
async def transcribir_audio(audio_bytes, extension="ogg"):
    """Transcribe audio usando Whisper en Groq."""
    loop = asyncio.get_event_loop()

    def llamar():
        with tempfile.NamedTemporaryFile(suffix=f".{extension}", delete=False) as tmp:
            tmp.write(audio_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, "rb") as f:
                return client.audio.transcriptions.create(
                    model=MODELO_VOZ,
                    file=f,
                    language="es"
                )
        finally:
            os.unlink(tmp_path)

    try:
        resp = await loop.run_in_executor(None, llamar)
        return resp.text
    except Exception as e:
        logger.error(f"Error transcribiendo audio: {e}")
        return None



# ─────────────────────────────────────────────
# ONBOARDING — Conversación de bienvenida
# ─────────────────────────────────────────────
async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)

    if perfil.get("onboarding") == 1:
        nombre = perfil.get("nombre", "")
        await update.message.reply_text(
            f"¡Hola de nuevo{', ' + nombre if nombre else ''}! ¿En qué te ayudo hoy?",
            parse_mode="Markdown"
        )
        return ConversationHandler.END

    await update.message.reply_text(
        "¡Hola! Soy *JARVIS*, tu asistente académico.\n\n"
        "Antes de empezar, necesito conocerte. ¿Cómo te llamas?",
        parse_mode="Markdown"
    )
    return OB_NOMBRE

async def ob_recibir_nombre(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    texto   = update.message.text.strip()

    # Extraer nombre real — ignorar frases introductorias comunes
    PREFIJOS = [
        "me llamo", "mi nombre es", "soy", "me dicen",
        "me llaman", "llamame", "llámame", "llámame",
        "mi nombre:", "nombre:", "hola soy", "hola me llamo"
    ]
    texto_lower = texto.lower()
    nombre_raw  = texto
    for pref in PREFIJOS:
        if texto_lower.startswith(pref):
            nombre_raw = texto[len(pref):].strip()
            break

    # Tomar solo el primer nombre (no apellidos completos)
    nombre = nombre_raw.split()[0].capitalize() if nombre_raw.split() else texto.split()[0].capitalize()

    context.user_data["ob_nombre"] = nombre
    await update.message.reply_text(
        f"Buena, {nombre}. ¿Qué carrera estudias?\n"
        f"_(Ej: Ingeniería Civil, Sistemas, Mecánica, etc.)_",
        parse_mode="Markdown"
    )
    return OB_CARRERA

async def ob_recibir_carrera(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["ob_carrera"] = update.message.text.strip().capitalize()
    await update.message.reply_text(
        "¿En qué ciclo vas? _(Ej: 3, 5to, Séptimo)_",
        parse_mode="Markdown"
    )
    return OB_CICLO

async def ob_recibir_ciclo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["ob_ciclo"] = update.message.text.strip()
    await update.message.reply_text(
        "¿En qué universidad estudias?",
        parse_mode="Markdown"
    )
    return OB_UNIVERSIDAD

async def ob_recibir_universidad(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id  = update.effective_user.id
    perfil   = cargar_perfil(user_id)
    perfil.update({
        "nombre":      context.user_data.get("ob_nombre"),
        "carrera":     context.user_data.get("ob_carrera"),
        "ciclo":       context.user_data.get("ob_ciclo"),
        "universidad": update.message.text.strip(),
        "onboarding":  1
    })
    guardar_perfil(user_id, perfil)

    nombre    = perfil["nombre"]
    carrera   = perfil["carrera"]
    ciclo     = perfil["ciclo"]
    uni       = perfil["universidad"]

    await update.message.reply_text(
        f"Listo, {nombre}. Todo guardado:\n\n"
        f"• Carrera: {carrera}\n"
        f"• Ciclo: {ciclo}\n"
        f"• Universidad: {uni}\n\n"
        f"Puedo ayudarte con ejercicios, PDFs, Excel, presentaciones, "
        f"recordatorios y más. ¿Por dónde empezamos?",
        parse_mode="Markdown"
    )
    return ConversationHandler.END

async def ob_cancelar(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)
    perfil["onboarding"] = 1
    guardar_perfil(user_id, perfil)
    await update.message.reply_text("Ok, empecemos. ¿Qué necesitas?")
    return ConversationHandler.END


# ─────────────────────────────────────────────
# COMANDOS
# ─────────────────────────────────────────────
async def cmd_ayuda(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "*JARVIS 3.0 — Lo que puedo hacer:*\n\n"
        "📐 *Ejercicios* — Física, Cálculo, Química, Ingeniería\n"
        "📄 *PDFs* — Soluciones con fórmulas renderizadas\n"
        "📊 *Excel* — Tablas, datos, fórmulas automáticas\n"
        "📑 *Presentaciones* — PowerPoint con diseño profesional\n"
        "📝 *Word* — Informes y documentos formateados\n"
        "🎤 *Voz* — Mándame un audio, lo entiendo\n"
        "📸 *Imágenes* — Foto de tu ejercicio, lo resuelvo\n"
        "🔔 *Recordatorios* — \"Recuérdame el examen el viernes a las 3pm\"\n"
        "💰 *Contabilidad* — \"Gasté 15 soles en el comedor\"\n\n"
        "*Comandos:*\n"
        "/perfil — Ver tu información\n"
        "/recordatorios — Ver tus recordatorios\n"
        "/contabilidad — Ver tu balance\n"
        "/limpiar — Borrar historial de chat\n"
        "/yo — Actualizar tu perfil",
        parse_mode="Markdown"
    )

async def cmd_perfil(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)
    now     = datetime.now()
    conn    = sqlite3.connect(DB_PATH)
    c       = conn.cursor()
    c.execute("SELECT mensajes, plan FROM uso_mensual WHERE user_id=? AND anio=? AND mes=?",
              (user_id, now.year, now.month))
    row = c.fetchone()
    conn.close()
    mensajes, plan = (row or (0, "basico"))
    limite = PLAN_FULL_LIMITE if plan == "full" else PLAN_BASICO_LIMITE

    await update.message.reply_text(
        f"*Tu perfil:*\n\n"
        f"👤 Nombre: {perfil.get('nombre', 'No registrado')}\n"
        f"🎓 Carrera: {perfil.get('carrera', 'No registrada')}\n"
        f"📚 Ciclo: {perfil.get('ciclo', 'No registrado')}\n"
        f"🏛 Universidad: {perfil.get('universidad', 'No registrada')}\n\n"
        f"📊 Plan: {plan.upper()}\n"
        f"💬 Mensajes este mes: {mensajes}/{limite}",
        parse_mode="Markdown"
    )

async def cmd_yo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Reinicia el onboarding para actualizar perfil."""
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)
    perfil["onboarding"] = 0
    guardar_perfil(user_id, perfil)
    await update.message.reply_text(
        "Vamos a actualizar tu perfil. ¿Cómo te llamas?"
    )
    return OB_NOMBRE

async def cmd_limpiar(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    guardar_historial(user_id, [])
    guardar_resumen(user_id, "")
    await update.message.reply_text("Historial limpiado. Empezamos de cero.")

async def cmd_recordatorios(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    conn    = sqlite3.connect(DB_PATH)
    c       = conn.cursor()
    c.execute("SELECT titulo, fecha_hora FROM recordatorios WHERE user_id=? AND enviado=0 ORDER BY fecha_hora",
              (user_id,))
    rows = c.fetchall()
    conn.close()

    if not rows:
        await update.message.reply_text("No tienes recordatorios pendientes.")
        return

    texto = "*Tus recordatorios pendientes:*\n\n"
    for titulo, fh in rows:
        try:
            dt = datetime.fromisoformat(fh)
            fecha_str = dt.strftime("%a %d/%m — %H:%M")
        except:
            fecha_str = fh
        texto += f"🔔 *{titulo}*\n   📅 {fecha_str}\n\n"

    await update.message.reply_text(texto, parse_mode="Markdown")

async def cmd_contabilidad(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Muestra resumen financiero completo usando módulo finanzas.py"""
    user_id = update.effective_user.id
    texto = resumen_mensual(user_id)
    await update.message.reply_text(texto, parse_mode="Markdown")


async def cmd_horario(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Muestra resumen del día académico."""
    user_id = update.effective_user.id

    if necesita_horario(user_id):
        context.user_data["esperando_foto"] = FOTO_HORARIO
        await update.message.reply_text(mensaje_pedir_foto_horario())
        return

    resumen = resumen_hoy(user_id)
    await update.message.reply_text(resumen, parse_mode="Markdown")


async def cmd_finanzas(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Alias de /contabilidad más amigable."""
    await cmd_contabilidad(update, context)


# ─────────────────────────────────────────────
# ENVÍO DE RECORDATORIO (scheduler)
# ─────────────────────────────────────────────
async def enviar_recordatorio(app, user_id, rec_id, titulo, mensaje):
    try:
        await app.bot.send_message(
            chat_id=user_id,
            text=f"🔔 *RECORDATORIO*\n\n{titulo}\n\n{mensaje}",
            parse_mode="Markdown"
        )
        conn = sqlite3.connect(DB_PATH)
        conn.execute("UPDATE recordatorios SET enviado=1 WHERE id=?", (rec_id,))
        conn.commit()
        conn.close()
    except Exception as e:
        logger.error(f"Error enviando recordatorio {rec_id}: {e}")


# ─────────────────────────────────────────────
# HELPERS DE MENSAJES (scope global)
# ─────────────────────────────────────────────
async def safe_edit(msg, texto_nuevo):
    try:
        await msg.edit_text(texto_nuevo)
    except Exception:
        pass

async def safe_delete(msg):
    try:
        await msg.delete()
    except Exception:
        pass


# ─────────────────────────────────────────────
# PROCESADOR CENTRAL DE TEXTO
# ─────────────────────────────────────────────
async def procesar_texto(update: Update, context: ContextTypes.DEFAULT_TYPE, texto: str):
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)

    # Registrar usuario en tabla central
    registrar_usuario(user_id, update.effective_user.username or "")

    # Verificar suscripción activa
    ok, n_msgs = verificar_limite(user_id, perfil)
    if not ok:
        if not suscripcion_activa(user_id):
            # ── Mensaje al alumno ──
            user  = update.effective_user
            uname = f"@{user.username}" if user.username else user.first_name or "Sin nombre"
            await update.message.reply_text(
                "🔒 *Acceso restringido*\n\n"
                "JARVIS es un servicio de pago para estudiantes.\n\n"
                "Para activar tu acceso:\n"
                "👉 Escríbele a @uziel0509 con tu nombre y carrera\n"
                "💰 Plan mensual: S/ 30/mes\n\n"
                "Tu tutor IA disponible 24/7 para resolver cualquier ejercicio.",
                parse_mode="Markdown"
            )
            # ── Notificación al admin ──
            if ADMIN_ID:
                try:
                    await context.bot.send_message(
                        chat_id=ADMIN_ID,
                        text=(
                            f"🔔 *Alumno sin acceso intentó usar el bot*\n\n"
                            f"👤 Nombre: {user.first_name or 'Sin nombre'} {user.last_name or ''}\n"
                            f"🆔 User ID: `{user_id}`\n"
                            f"📱 Username: {uname}\n"
                            f"💬 Mensaje: _{texto[:100]}_\n\n"
                            f"Para activarlo 1 mes:\n"
                            f"`/activar {user_id} 1`"
                        ),
                        parse_mode="Markdown"
                    )
                except Exception as e:
                    logger.warning(f"No se pudo notificar al admin: {e}")
        else:
            plan = perfil.get("plan", "basico")
            limite = PLAN_BASICO_LIMITE if plan == "basico" else PLAN_FULL_LIMITE
            await update.message.reply_text(
                f"Llegaste al límite de {limite} mensajes del plan {plan.upper()}.\n\n"
                f"Para más mensajes escríbeme a @uziel0509."
            )
        return

    # Detectar intención

    # ── SALUDO CORTO — respuesta directa sin modelo ──
    _t = texto.lower().strip()
    _saludos = {"hola","hello","hi","ola","buenas","hey","que tal","qué tal",
                "buenas tardes","buenas noches","buenos dias","buenos días",
                "buen dia","buen día"}
    if _t in _saludos:
        _nombre = perfil.get("nombre", "")
        await update.message.reply_text(f"Hola{chr(32) + _nombre if _nombre else chr(46)} ¿Qué necesitas?")
        incrementar_uso(user_id)
        return

    intencion = detectar_intencion(texto)
    modelo    = elegir_modelo(intencion)

    # ── RECORDATORIO ──
    if intencion == "recordatorio":
        dt = parsear_fecha_natural(texto)
        titulo_rec = extraer_titulo_recordatorio(texto)
        if dt:
            conn = sqlite3.connect(DB_PATH)
            c    = conn.cursor()
            c.execute(
                "INSERT INTO recordatorios (user_id, titulo, fecha_hora, mensaje) VALUES (?,?,?,?)",
                (user_id, titulo_rec, dt.isoformat(), titulo_rec)
            )
            rec_id = c.lastrowid
            conn.commit()
            conn.close()

            # Programar en scheduler
            context.application.job_queue.run_once(
                lambda ctx: asyncio.create_task(
                    enviar_recordatorio(context.application, user_id, rec_id, titulo_rec, titulo_rec)
                ),
                when=dt,
                name=f"rec_{rec_id}"
            )

            await update.message.reply_text(
                f"✅ *Recordatorio guardado*\n\n"
                f"📌 {titulo_rec}\n"
                f"📅 {dt.strftime('%A %d/%m/%Y a las %H:%M')}",
                parse_mode="Markdown"
            )
        else:
            await update.message.reply_text(
                "No entendí la fecha/hora. Dime algo como:\n"
                "\"Recuérdame el examen el viernes a las 3pm\""
            )
        incrementar_uso(user_id)
        return

    # ── HORARIO / AGENDA ──
    if intencion == "horario":
        if necesita_horario(user_id):
            context.user_data["esperando_foto"] = FOTO_HORARIO
            await update.message.reply_text(mensaje_pedir_foto_horario())
        else:
            await update.message.reply_text(resumen_hoy(user_id), parse_mode="Markdown")
        incrementar_uso(user_id)
        return

    # ── REGISTRAR EXAMEN ──
    if intencion == "examen":
        dt = parsear_fecha_natural(texto)
        t_lower = texto.lower()
        materia = "Materia"
        for kw in ["examen de ", "examen del ", "parcial de ", "final de "]:
            if kw in t_lower:
                resto  = t_lower.split(kw, 1)[1]
                materia = resto.split(" el ")[0].split(" mañana")[0].strip().capitalize()
                break
        if dt:
            registrar_examen(user_id, materia, dt.strftime("%Y-%m-%d"), dt.strftime("%H:%M"))
            dt_aviso = dt - timedelta(days=1)
            if dt_aviso > datetime.now():
                context.application.job_queue.run_once(
                    lambda ctx, uid=user_id, m=materia, d=dt: asyncio.create_task(
                        context.application.bot.send_message(
                            chat_id=uid,
                            text=f"📚 *Recordatorio*: Mañana tienes examen de *{m}* a las {d.strftime('%H:%M')}. ¡Repasa esta noche!",
                            parse_mode="Markdown"
                        )
                    ),
                    when=dt_aviso,
                    name=f"exam_aviso_{user_id}_{materia}"
                )
            await update.message.reply_text(
                f"📝 *Examen registrado*\n\n"
                f"📖 Materia: {materia}\n"
                f"📅 Fecha: {dt.strftime('%A %d/%m/%Y a las %H:%M')}\n\n"
                f"Te avisaré un día antes. 💪",
                parse_mode="Markdown"
            )
        else:
            await update.message.reply_text(
                "No entendí la fecha. Dime algo como:\n"
                "\"Tengo examen de Física el viernes a las 10am\""
            )
        incrementar_uso(user_id)
        return

    # ── REGISTRAR ENTREGA ──
    if intencion == "entrega":
        dt = parsear_fecha_natural(texto)
        t_lower  = texto.lower()
        concepto = "Entrega"
        for kw in ["entrega de ", "entregar ", "presentar "]:
            if kw in t_lower:
                resto    = t_lower.split(kw, 1)[1]
                concepto = resto.split(" el ")[0].split(" mañana")[0].strip().capitalize()
                break
        if dt:
            registrar_entrega(user_id, concepto, dt.strftime("%Y-%m-%d"), dt.strftime("%H:%M"))
            await update.message.reply_text(
                f"✅ *Entrega registrada*\n\n"
                f"📋 {concepto}\n"
                f"📅 {dt.strftime('%A %d/%m/%Y a las %H:%M')}",
                parse_mode="Markdown"
            )
        else:
            await update.message.reply_text(
                "No entendí la fecha. Dime algo como:\n"
                "\"Tengo que entregar el informe el jueves a las 8am\""
            )
        incrementar_uso(user_id)
        return

    # ── CONTABILIDAD ──
    if intencion == "contabilidad":
        t_lower = texto.lower()
        # Ver balance / resumen
        if any(k in t_lower for k in ["mi balance", "mis gastos", "cuánto llevo", "cuanto llevo",
                                        "mi saldo", "mis ingresos", "ver gastos", "resumen"]):
            await update.message.reply_text(resumen_mensual(user_id), parse_mode="Markdown")
        # Configurar límite mensual
        elif any(k in t_lower for k in ["limite", "límite", "presupuesto"]):
            m = re.search(r'(\d+(?:\.\d{1,2})?)', texto)
            if m:
                monto_limite = float(m.group(1))
                configurar_limite(user_id, monto_limite)
                await update.message.reply_text(
                    f"✅ Límite mensual configurado en *S/ {monto_limite:.2f}*\n"
                    f"Te avisaré cuando llegues al 80% y al 100%.",
                    parse_mode="Markdown"
                )
            else:
                await update.message.reply_text("¿Cuánto quieres de límite mensual? Ej: \"Límite 800 soles\"")
        else:
            # Registrar movimiento usando módulo finanzas.py
            tipo, monto, concepto = parsear_movimiento(texto)
            if monto > 0:
                if tipo == "gasto":
                    alerta = registrar_gasto(user_id, monto, concepto)
                else:
                    alerta = registrar_ingreso(user_id, monto, concepto)
                icono = "💚" if tipo == "ingreso" else "🔴"
                respuesta = (
                    f"{icono} Registrado: *{concepto}*\n"
                    f"Monto: S/ {monto:.2f} ({tipo})\n\n"
                    f"Usa /finanzas para ver tu balance."
                )
                if alerta:
                    respuesta += f"\n\n{alerta}"
                await update.message.reply_text(respuesta, parse_mode="Markdown")
            else:
                await update.message.reply_text(
                    "No detecté el monto. Dime algo como:\n"
                    "\"Gasté 25 soles en libros\""
                )
        incrementar_uso(user_id)
        return

    # ── GENERACIÓN DE ARCHIVOS ──
    if intencion in ("excel", "presentacion", "documento"):
        msg_espera = await update.message.reply_text(
            "⏳ Generando tu archivo..." if intencion != "presentacion"
            else "⏳ Creando presentación, dame un momento..."
        )
        estructura = await generar_estructura_archivo(intencion, texto, perfil)
        if not estructura:
            try:
                await msg_espera.edit_text("No pude generar el archivo. Intenta describir mejor lo que necesitas.")
            except Exception:
                pass
            return

        try:
            if intencion == "excel":
                path = crear_excel(estructura, user_id, estructura.get("titulo", "Datos"))
                try:
                    await msg_espera.delete()
                except Exception:
                    pass
                await update.message.reply_document(
                    document=open(path, "rb"),
                    filename=f"{estructura.get('titulo', 'datos')}.xlsx",
                    caption=f"📊 *{estructura.get('titulo', 'Excel')}*\nGenerado por JARVIS",
                    parse_mode="Markdown"
                )
            elif intencion == "presentacion":
                path = crear_presentacion(
                    estructura.get("slides", []),
                    user_id,
                    estructura.get("titulo", "Presentación")
                )
                try:
                    await msg_espera.delete()
                except Exception:
                    pass
                await update.message.reply_document(
                    document=open(path, "rb"),
                    filename=f"{estructura.get('titulo', 'presentacion')}.pptx",
                    caption=f"📑 *{estructura.get('titulo', 'Presentación')}*\nGenerado por JARVIS",
                    parse_mode="Markdown"
                )
            elif intencion == "documento":
                path = crear_documento_word(
                    estructura.get("contenido", ""),
                    user_id,
                    estructura.get("titulo", "Documento")
                )
                try:
                    await msg_espera.delete()
                except Exception:
                    pass
                await update.message.reply_document(
                    document=open(path, "rb"),
                    filename=f"{estructura.get('titulo', 'documento')}.docx",
                    caption=f"📝 *{estructura.get('titulo', 'Documento')}*\nGenerado por JARVIS",
                    parse_mode="Markdown"
                )
        except Exception as e:
            logger.error(f"Error generando archivo {intencion}: {e}")
            try:
                await msg_espera.edit_text(
                    "Hubo un error generando el archivo. Intenta con una descripción más específica."
                )
            except Exception:
                try:
                    await update.message.reply_text("Tuve un problema, intenta de nuevo.")
                except Exception:
                    pass
        incrementar_uso(user_id)
        return

    # ── CHAT / EJERCICIO ──
    msg_espera = await update.message.reply_text("⏳")

    try:
        loop  = asyncio.get_event_loop()
        msgs, historial = construir_payload(user_id, texto, perfil)

        def llamar_ia():
            # Fallback automático: si el modelo principal falla → usa MODELO_CHAT
            modelos_intentar = [modelo]
            if modelo != MODELO_CHAT:
                modelos_intentar.append(MODELO_CHAT)
            last_error = None
            for m in modelos_intentar:
                try:
                    return client.chat.completions.create(
                        model=m,
                        messages=msgs,
                        max_tokens=4000
                    )
                except Exception as e:
                    last_error = e
                    logger.warning(f"Modelo {m} falló: {e}. Intentando fallback...")
            raise last_error

        resp = await loop.run_in_executor(None, llamar_ia)
        respuesta = resp.choices[0].message.content

        # Actualizar historial
        historial.append({"role": "user", "content": texto})
        historial.append({"role": "assistant", "content": respuesta})
        guardar_historial(user_id, historial)

        # Generar resumen si hace falta
        if len(historial) >= SUMMARY_EVERY:
            asyncio.create_task(generar_resumen(user_id, historial))

        incrementar_uso(user_id)

        # Decidir si generar PDF automático
        es_ejercicio    = intencion in ("ejercicio_dificil", "ejercicio_simple")
        respuesta_larga = len(respuesta) > 600
        tiene_pasos     = bool(re.search(r'paso\s+\d+|step\s+\d+|\d+\.', respuesta, re.IGNORECASE))

        if es_ejercicio and (respuesta_larga or tiene_pasos):
            await safe_delete(msg_espera)
            try:
                from modulos.agente_academico import resolver_y_generar_pdf
                titulo_pdf = f"Ejercicio — {datetime.now().strftime('%d/%m/%Y')}"
                path_pdf, _ = await loop.run_in_executor(
                    None,
                    lambda: resolver_y_generar_pdf(client, texto, user_id, titulo_pdf, perfil)
                )
                await update.message.reply_document(
                    document=open(path_pdf, "rb"),
                    filename="solucion_jarvis.pdf",
                    caption="📄 Aquí tienes tu solución completa"
                )
            except Exception as e:
                logger.error(f"Error generando PDF: {e}")
                elementos      = procesar_output(respuesta)
                respuesta_chat = elementos_a_texto_plano(elementos)
                await update.message.reply_text(respuesta_chat[:8000])
        else:
            await safe_edit(msg_espera, limpiar_latex(respuesta))

    except Exception as e:
        logger.error(f"Error en procesar_texto: {e}")
        await safe_edit(msg_espera,
            "Tuve un problema procesando tu mensaje. Intenta de nuevo en un momento."
        )


# ─────────────────────────────────────────────
# HANDLERS PRINCIPALES
# ─────────────────────────────────────────────
async def manejar_mensaje(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.message and update.message.text:
        await procesar_texto(update, context, update.message.text)

async def manejar_imagen(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)

    ok, _ = verificar_limite(user_id, perfil)
    if not ok:
        await update.message.reply_text("Límite mensual alcanzado. Actualiza tu plan para continuar.")
        return

    # ── Detectar tipo de foto esperada por contexto ──
    tipo_foto = context.user_data.pop("esperando_foto", None)

    # Si no hay en context.user_data, revisar archivo temporal (cron post-examen)
    if tipo_foto is None:
        estado_path = Path(f"{PERFILES_DIR}/{user_id}/esperando_examen.json")
        if estado_path.exists():
            try:
                estado = json.loads(estado_path.read_text())
                # Solo válido si es del día de hoy
                ts = datetime.fromisoformat(estado["ts"])
                if (datetime.now() - ts).total_seconds() < 86400:  # 24h
                    tipo_foto = FOTO_EXAMEN
                    context.user_data["examen_id"]     = estado.get("exam_id")
                    context.user_data["examen_materia"] = estado.get("materia", "la materia")
                    estado_path.unlink()  # borrar el archivo después de leerlo
            except Exception:
                pass
    if tipo_foto is None:
        tipo_foto = FOTO_EJERCICIO

    msg_espera = await update.message.reply_text(
        "📅 Procesando tu horario..." if tipo_foto == FOTO_HORARIO else
        "📝 Analizando tu examen..." if tipo_foto == FOTO_EXAMEN else
        "🔍 Analizando imagen..."
    )

    try:
        photo   = update.message.photo[-1]
        caption = (update.message.caption or "").strip()
        loop    = asyncio.get_event_loop()

        # Descargar imagen
        file_obj  = await context.bot.get_file(photo.file_id)
        img_bytes = await loop.run_in_executor(
            None, lambda: __import__('requests').get(file_obj.file_path).content
        )
        img_b64 = base64.b64encode(img_bytes).decode("utf-8")

        # ── FOTO HORARIO ──
        if tipo_foto == FOTO_HORARIO:
            def extraer_horario():
                return client.chat.completions.create(
                    model=MODELO_VISION,
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                            {"type": "text", "text": PROMPT_EXTRAER_HORARIO}
                        ]
                    }],
                    max_tokens=2000
                )
            resp_vision = await loop.run_in_executor(None, extraer_horario)
            raw = resp_vision.choices[0].message.content
            try:
                import re as _re
                json_match = _re.search(r'\[[\s\S]*\]|\{[\s\S]*\}', raw)
                if json_match:
                    parsed = json.loads(json_match.group())
                    # Normalizar: puede venir como lista de clases o como dict con key "clases"
                    if isinstance(parsed, dict):
                        clases = parsed.get("clases", [])
                    else:
                        clases = parsed
                else:
                    clases = []

                if clases:
                    guardar_horario_extraido(user_id, clases)
                    await safe_delete(msg_espera)
                    await update.message.reply_text(
                        f"✅ *Horario cargado exitosamente*\n\n"
                        f"📚 {len(clases)} clases detectadas\n\n"
                        f"Usa /horario para ver tu agenda del día.",
                        parse_mode="Markdown"
                    )
                else:
                    await safe_delete(msg_espera)
                    await update.message.reply_text(
                        "No pude leer el horario claramente. Intenta con una foto más nítida."
                    )
            except Exception as e:
                logger.error(f"Error parsing horario: {e}")
                await safe_delete(msg_espera)
                await update.message.reply_text("Error procesando el horario. Intenta de nuevo.")
            incrementar_uso(user_id)
            return

        # ── FOTO EXAMEN (post-examen para análisis de errores) ──
        if tipo_foto == FOTO_EXAMEN:
            materia = context.user_data.pop("examen_materia", "la materia")
            examen_id = context.user_data.pop("examen_id", None)

            def analizar_examen():
                return client.chat.completions.create(
                    model=MODELO_EJERCICIOS,
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                            {"type": "text", "text": PROMPT_ANALIZAR_EXAMEN.format(materia=materia)}
                        ]
                    }],
                    max_tokens=3000
                )
            resp_analisis = await loop.run_in_executor(None, analizar_examen)
            analisis = resp_analisis.choices[0].message.content

            if examen_id:
                guardar_resultado_examen(user_id, examen_id, materia, analisis)

            await safe_delete(msg_espera)
            await update.message.reply_text(limpiar_latex(analisis)[:4000])

            # Generar PDF de retroalimentación
            try:
                titulo_pdf = f"Análisis Examen {materia} — {datetime.now().strftime('%d/%m/%Y')}"
                path_pdf = _crear_pdf_limpio(analisis, user_id, titulo_pdf, perfil)
                await update.message.reply_document(
                    document=open(path_pdf, "rb"),
                    filename=f"analisis_{materia.lower().replace(' ','_')}.pdf",
                    caption=f"📊 Retroalimentación completa de tu examen de {materia}"
                )
            except Exception as e:
                logger.error(f"Error PDF análisis examen: {e}")
            incrementar_uso(user_id)
            return

        # ── FOTO EJERCICIO (flujo original) ──
        pide_pdf = bool(re.search(r'\bpdf\b|en pdf|como pdf|formato pdf|en documento', caption, re.IGNORECASE))

        contexto_visual, solucion, modelo_usado = await analizar_imagen_completo(
            img_bytes, caption, user_id, perfil
        )

        historial = cargar_historial(user_id)
        historial.append({"role": "user",      "content": f"[IMAGEN con instrucción: '{caption}'] {contexto_visual[:300]}"})
        historial.append({"role": "assistant", "content": solucion})
        guardar_historial(user_id, historial)
        incrementar_uso(user_id)

        # Usar pre_render para limpiar la respuesta
        elementos     = procesar_output(solucion)
        solucion_chat = elementos_a_texto_plano(elementos)
        respuesta_larga = len(solucion) > 600
        tiene_pasos     = bool(re.search(r'paso\s+\d+|\d+[\.\)]\s', solucion, re.IGNORECASE))

        await safe_delete(msg_espera)

        if not pide_pdf:
            await update.message.reply_text(solucion_chat[:8000])

        if pide_pdf or respuesta_larga or tiene_pasos:
            try:
                from modulos.agente_academico import resolver_y_generar_pdf
                titulo_pdf = f"Ejercicio — {datetime.now().strftime('%d/%m/%Y')}"
                path_pdf, _ = await loop.run_in_executor(
                    None,
                    lambda: resolver_y_generar_pdf(client, solucion, user_id, titulo_pdf, perfil)
                )
                await update.message.reply_document(
                    document=open(path_pdf, "rb"),
                    filename="solucion_jarvis.pdf",
                    caption="📄 Solución completa en PDF"
                )
            except Exception as e:
                logger.error(f"Error PDF imagen: {e}")

    except Exception as e:
        logger.error(f"Error manejar_imagen ({tipo_foto}): {e}")
        try:
            await msg_espera.edit_text("No pude procesar la imagen. Intenta de nuevo.")
        except Exception:
            pass

    # ── BLOQUE VACÍO — mantener compatibilidad con código previo ──
    # (el resto del handler original está integrado arriba)
    return  # evita caer en código legacy

async def manejar_voz(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    perfil  = cargar_perfil(user_id)

    ok, _ = verificar_limite(user_id, perfil)
    if not ok:
        await update.message.reply_text("Límite mensual alcanzado.")
        return

    msg_espera = await update.message.reply_text("🎤 Transcribiendo audio...")

    try:
        voice = update.message.voice or update.message.audio
        loop  = asyncio.get_event_loop()
        file_obj   = await context.bot.get_file(voice.file_id)
        audio_bytes = await loop.run_in_executor(
            None, lambda: __import__('requests').get(file_obj.file_path).content
        )

        texto_transcrito = await transcribir_audio(audio_bytes, "ogg")

        if not texto_transcrito:
            try:
                await msg_espera.edit_text("No pude entender el audio. Intenta hablar más claro.")
            except Exception:
                pass
            return

        try:
            await msg_espera.edit_text(f"🎤 _{texto_transcrito}_", parse_mode="Markdown")
        except Exception:
            pass
        await procesar_texto(update, context, texto_transcrito)

    except Exception as e:
        logger.error(f"Error manejar_voz: {e}")
        try:
            await msg_espera.edit_text("Error procesando el audio.")
        except Exception:
            pass


# ─────────────────────────────────────────────
# RECARGA DE RECORDATORIOS AL INICIAR
# ─────────────────────────────────────────────
async def cmd_finanzas(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Muestra resumen financiero mensual completo."""
    user_id = update.effective_user.id
    await update.message.reply_text(resumen_mensual(user_id), parse_mode="Markdown")


async def cmd_horario(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Muestra la agenda del día."""
    user_id = update.effective_user.id
    if necesita_horario(user_id):
        context.user_data["esperando_foto"] = FOTO_HORARIO
        await update.message.reply_text(mensaje_pedir_foto_horario())
    else:
        await update.message.reply_text(resumen_hoy(user_id), parse_mode="Markdown")



# ─────────────────────────────────────────────
# COMANDOS DE ADMINISTRACIÓN (solo admin)
# ─────────────────────────────────────────────

async def cmd_activar(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    Activa suscripción de un alumno.
    Uso: /activar <user_id> <meses> [plan]
    Ejemplo: /activar 123456789 1 full
    """
    uid = update.effective_user.id
    if uid != ADMIN_ID:
        await update.message.reply_text("⛔ Sin acceso.")
        return

    args = context.args
    if not args or len(args) < 2:
        await update.message.reply_text(
            "Uso: /activar <user_id> <meses> [plan]\n"
            "Ejemplo: /activar 123456789 1 full"
        )
        return

    try:
        target_id = int(args[0])
        meses     = int(args[1])
        plan      = args[2] if len(args) > 2 else "full"
        venc      = activar_suscripcion(target_id, plan, meses, uid)
        await update.message.reply_text(
            f"✅ *Suscripción activada*\n\n"
            f"👤 User ID: `{target_id}`\n"
            f"📦 Plan: {plan.upper()}\n"
            f"📅 Vence: {venc}",
            parse_mode="Markdown"
        )
        # Notificar al alumno
        try:
            await context.bot.send_message(
                chat_id=target_id,
                text=f"🎉 *¡Tu acceso a JARVIS está activado!*\n\n"
                     f"Plan: {plan.upper()} — Vence: {venc}\n\n"
                     f"Ya puedes empezar. Escribe lo que necesitas o usa /ayuda.",
                parse_mode="Markdown"
            )
        except Exception:
            await update.message.reply_text("(No se pudo notificar al alumno directamente)")
    except (ValueError, IndexError):
        await update.message.reply_text("Formato incorrecto. Usa: /activar 123456789 1")


async def cmd_revocar(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Revoca suscripción. Uso: /revocar <user_id>"""
    uid = update.effective_user.id
    if uid != ADMIN_ID:
        await update.message.reply_text("⛔ Sin acceso.")
        return

    args = context.args
    if not args:
        await update.message.reply_text("Uso: /revocar <user_id>")
        return

    try:
        target_id = int(args[0])
        conn = sqlite3.connect(DB_PATH)
        conn.execute(
            "UPDATE suscripciones SET plan='inactivo' WHERE user_id=?",
            (target_id,)
        )
        conn.commit()
        conn.close()
        await update.message.reply_text(f"✅ Suscripción revocada para `{target_id}`", parse_mode="Markdown")
    except ValueError:
        await update.message.reply_text("User ID inválido.")


async def cmd_alumnos(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Lista todos los alumnos con suscripción activa."""
    uid = update.effective_user.id
    if uid != ADMIN_ID:
        await update.message.reply_text("⛔ Sin acceso.")
        return

    conn = sqlite3.connect(DB_PATH)
    c    = conn.cursor()
    c.execute("""
        SELECT s.user_id, u.username, s.plan, s.fecha_vencimiento
        FROM suscripciones s
        LEFT JOIN usuarios u ON s.user_id = u.user_id
        WHERE s.plan != 'inactivo'
        ORDER BY s.fecha_vencimiento DESC
    """)
    rows = c.fetchall()
    conn.close()

    if not rows:
        await update.message.reply_text("No hay alumnos con suscripción activa.")
        return

    texto = "*📋 Alumnos activos*\n\n"
    ahora = datetime.now()
    for uid_al, username, plan, venc in rows:
        try:
            venc_dt = datetime.fromisoformat(venc) if venc else None
            dias    = (venc_dt - ahora).days if venc_dt else "∞"
            estado  = "✅" if (dias == "∞" or dias > 0) else "⚠️ VENCIDO"
        except Exception:
            dias, estado = "?", "⚠️"
        nombre_al = f"@{username}" if username else f"ID:{uid_al}"
        texto += f"{estado} {nombre_al} — {plan.upper()} — {dias}d restantes\n"

    await update.message.reply_text(texto, parse_mode="Markdown")


async def cmd_mystatus(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Muestra el estado de suscripción del usuario."""
    user_id = update.effective_user.id
    info = get_suscripcion(user_id)

    if not info["activo"]:
        await update.message.reply_text(
            "🔒 *Sin suscripción activa*\n\n"
            "Escríbele a @uziel0509para activar tu acceso.\n"
            "💰 Plan mensual: S/ 30/mes",
            parse_mode="Markdown"
        )
        return

    venc = info.get("vencimiento", "")
    try:
        venc_dt = datetime.fromisoformat(venc)
        dias    = (venc_dt - datetime.now()).days
        venc_str = venc_dt.strftime("%d/%m/%Y")
    except Exception:
        dias, venc_str = "∞", "Sin vencimiento"

    await update.message.reply_text(
        f"✅ *Suscripción activa*\n\n"
        f"📦 Plan: {info['plan'].upper()}\n"
        f"📅 Vence: {venc_str}\n"
        f"⏳ Días restantes: {dias}",
        parse_mode="Markdown"
    )

async def recargar_recordatorios_pendientes(app):
    conn = sqlite3.connect(DB_PATH)
    c    = conn.cursor()
    c.execute("SELECT id, user_id, titulo, fecha_hora, mensaje FROM recordatorios WHERE enviado=0")
    rows = c.fetchall()
    conn.close()
    count = 0
    ahora = datetime.now()
    for rec_id, user_id, titulo, fh, mensaje in rows:
        try:
            dt = datetime.fromisoformat(fh)
            if dt > ahora:
                app.job_queue.run_once(
                    lambda ctx, uid=user_id, rid=rec_id, t=titulo, m=mensaje:
                        asyncio.create_task(enviar_recordatorio(app, uid, rid, t, m)),
                    when=dt,
                    name=f"rec_{rec_id}"
                )
                count += 1
        except Exception as e:
            logger.warning(f"No se pudo reprogramar recordatorio {rec_id}: {e}")
    logger.info(f"Reprogramados {count} recordatorios pendientes")


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────
async def cron_pagos_recurrentes(context):
    """Cron diario: revisa y notifica pagos recurrentes del día."""
    from modulos.finanzas import get_pagos_hoy
    import sqlite3 as _sql
    conn = _sql.connect(DB_PATH)
    c    = conn.cursor()
    c.execute("SELECT DISTINCT user_id FROM uso_mensual")
    usuarios = [r[0] for r in c.fetchall()]
    conn.close()
    for uid in usuarios:
        pagos = get_pagos_hoy(uid)
        if pagos:
            lineas = "\n".join(f"• {p['concepto']}: S/ {p['monto']:.2f}" for p in pagos)
            try:
                await context.bot.send_message(
                    chat_id=uid,
                    text=f"💳 *Pagos recurrentes de hoy*\n\n{lineas}\n\n"
                         f"Recuerda registrarlos con /finanzas",
                    parse_mode="Markdown"
                )
            except Exception as e:
                logger.warning(f"No se pudo notificar pagos a {uid}: {e}")


async def cron_resumen_diario(context):
    """Cron cada mañana: envía resumen del día académico."""
    import sqlite3 as _sql
    conn = _sql.connect(DB_PATH)
    c    = conn.cursor()
    c.execute("SELECT DISTINCT user_id FROM uso_mensual")
    usuarios = [r[0] for r in c.fetchall()]
    conn.close()
    for uid in usuarios:
        try:
            from modulos.horario import necesita_horario, resumen_hoy
            if not necesita_horario(uid):
                resumen = resumen_hoy(uid)
                if "📚" in resumen or "📝" in resumen:  # Solo si hay algo que reportar
                    await context.bot.send_message(
                        chat_id=uid,
                        text=f"🌅 *Buenos días* — Aquí tu agenda de hoy:\n\n{resumen}",
                        parse_mode="Markdown"
                    )
        except Exception as e:
            logger.warning(f"No se pudo enviar resumen diario a {uid}: {e}")


async def cron_detectar_examenes_pasados(context):
    """Cron post-examen: detecta exámenes que ya pasaron y pide foto al alumno."""
    import sqlite3 as _sql
    ahora = datetime.now()
    conn = _sql.connect(DB_PATH)
    c    = conn.cursor()
    # Buscar exámenes de las últimas 6h que no han sido analizados
    c.execute("""
        SELECT id, user_id, materia, fecha, hora
        FROM examenes_pendientes
        WHERE analizado=0 AND notificado=0
        AND datetime(fecha || ' ' || hora) < datetime('now', 'localtime')
        AND datetime(fecha || ' ' || hora) > datetime('now', '-6 hours', 'localtime')
    """)
    rows = c.fetchall()
    for exam_id, uid, materia, fecha, hora in rows:
        try:
            await context.bot.send_message(
                chat_id=uid,
                text=f"📋 *¿Cómo te fue en el examen de {materia}?*\n\n"
                     f"Mándame una foto de tu examen y te doy retroalimentación detallada:\n"
                     f"• Errores conceptuales vs de cálculo\n"
                     f"• Qué reforzar antes del siguiente\n"
                     f"• PDF con análisis completo",
                parse_mode="Markdown"
            )
            # Marcar como notificado — cuando mande la foto se procesa
            c.execute("UPDATE examenes_pendientes SET notificado=1 WHERE id=?", (exam_id,))
            conn.commit()
            # Guardar estado en archivo temporal para que manejar_imagen sepa qué esperar
            estado_path = Path(f"{PERFILES_DIR}/{uid}/esperando_examen.json")
            estado_path.parent.mkdir(parents=True, exist_ok=True)
            estado_path.write_text(json.dumps({
                "exam_id": exam_id,
                "materia": materia,
                "ts": datetime.now().isoformat()
            }))
        except Exception as e:
            logger.warning(f"No se pudo notificar examen pasado a {uid}: {e}")
    conn.close()


def main():
    init_db()

    app = Application.builder().token(TELEGRAM_TOKEN).build()

    # Onboarding handler
    onboarding_handler = ConversationHandler(
        entry_points=[CommandHandler("start", cmd_start)],
        states={
            OB_NOMBRE:      [MessageHandler(filters.TEXT & ~filters.COMMAND, ob_recibir_nombre)],
            OB_CARRERA:     [MessageHandler(filters.TEXT & ~filters.COMMAND, ob_recibir_carrera)],
            OB_CICLO:       [MessageHandler(filters.TEXT & ~filters.COMMAND, ob_recibir_ciclo)],
            OB_UNIVERSIDAD: [MessageHandler(filters.TEXT & ~filters.COMMAND, ob_recibir_universidad)],
        },
        fallbacks=[CommandHandler("cancelar", ob_cancelar)],
        allow_reentry=True
    )

    app.add_handler(onboarding_handler)
    app.add_handler(CommandHandler("ayuda",         cmd_ayuda))
    app.add_handler(CommandHandler("help",          cmd_ayuda))
    app.add_handler(CommandHandler("perfil",        cmd_perfil))
    app.add_handler(CommandHandler("yo",            cmd_yo))
    app.add_handler(CommandHandler("limpiar",       cmd_limpiar))
    app.add_handler(CommandHandler("recordatorios", cmd_recordatorios))
    app.add_handler(CommandHandler("contabilidad",  cmd_finanzas))
    app.add_handler(CommandHandler("finanzas",      cmd_finanzas))
    app.add_handler(CommandHandler("horario",       cmd_horario))
    app.add_handler(CommandHandler("mystatus",      cmd_mystatus))
    app.add_handler(CommandHandler("activar",       cmd_activar))
    app.add_handler(CommandHandler("revocar",       cmd_revocar))
    app.add_handler(CommandHandler("alumnos",       cmd_alumnos))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, manejar_mensaje))
    app.add_handler(MessageHandler(filters.PHOTO,                   manejar_imagen))
    app.add_handler(MessageHandler(filters.VOICE | filters.AUDIO,   manejar_voz))

    async def post_init(application):
        await recargar_recordatorios_pendientes(application)

        # ── CRON JOBS ──
        # Resumen diario a las 7am hora Perú
        application.job_queue.run_daily(
            cron_resumen_diario,
            time=__import__('datetime').time(7, 0, 0),
            name="resumen_diario"
        )
        # Pagos recurrentes a las 8am
        application.job_queue.run_daily(
            cron_pagos_recurrentes,
            time=__import__('datetime').time(8, 0, 0),
            name="pagos_recurrentes"
        )
        # Detectar exámenes pasados cada 2h
        application.job_queue.run_repeating(
            cron_detectar_examenes_pasados,
            interval=7200,
            first=60,
            name="detectar_examenes_pasados"
        )

        # Menú de comandos actualizado
        await application.bot.set_my_commands([
            BotCommand("start",          "Iniciar / Bienvenida"),
            BotCommand("ayuda",          "Ver todo lo que puedo hacer"),
            BotCommand("horario",        "Ver agenda del día"),
            BotCommand("finanzas",       "Ver mi balance mensual"),
            BotCommand("recordatorios",  "Ver recordatorios pendientes"),
            BotCommand("mystatus",       "Ver mi suscripción"),
            BotCommand("perfil",         "Ver tu perfil"),
            BotCommand("limpiar",        "Borrar historial de chat"),
        ])

    app.post_init = post_init

    logger.info("JARVIS 3.2 iniciando — pagos, admin, fallback modelos, logging archivo")
    app.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == "__main__":
    main()
